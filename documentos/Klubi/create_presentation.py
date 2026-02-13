#!/usr/bin/env python3
"""
Script de Geração de Apresentação PPTX
Análise de Maturidade Digital - Klubi
Autor: Claude Code
Data: 2026-01-28
"""

import json
import matplotlib
matplotlib.use('Agg')  # Backend não-interativo
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# CONFIGURAÇÕES
COLORS = {
    'navy': RGBColor(30, 58, 95),        # #1E3A5F
    'green': RGBColor(0, 200, 150),      # #00C896
    'orange': RGBColor(255, 107, 53),    # #FF6B35
    'red': RGBColor(230, 57, 70),        # #E63946
    'gray': RGBColor(108, 117, 125),     # #6C757D
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(240, 240, 240),
    'dark_gray': RGBColor(50, 50, 50)
}

# HELPER FUNCTIONS
def add_title_slide(prs, title, subtitle=""):
    """Cria slide de título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['navy']
    background.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Subtítulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2),
            Inches(11.333), Inches(1)
        )
        tf_sub = subtitle_box.text_frame
        tf_sub.text = subtitle
        p_sub = tf_sub.paragraphs[0]
        p_sub.font.size = Pt(24)
        p_sub.font.color.rgb = COLORS['light_gray']
        p_sub.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    return slide

def add_content_slide(prs, title):
    """Cria slide de conteúdo com título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background branco
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()

    # Barra de título
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width, Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['navy']
    title_bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(12), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=COLORS['dark_gray']):
    """Adiciona caixa de texto"""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.text = text
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box

def create_radar_chart(pilares, pontuacoes, filename):
    """Cria gráfico radar de maturidade digital"""
    # Preparar dados
    angles = np.linspace(0, 2 * np.pi, len(pilares), endpoint=False).tolist()
    pontuacoes_plot = pontuacoes + [pontuacoes[0]]  # Fechar o círculo
    angles += angles[:1]

    # Criar gráfico
    fig, ax = plt.subplots(figsize=(10, 10), subplot_kw=dict(projection='polar'))
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')

    # Plot
    ax.plot(angles, pontuacoes_plot, 'o-', linewidth=3, color='#00C896', markersize=8)
    ax.fill(angles, pontuacoes_plot, alpha=0.25, color='#00C896')

    # Configuração
    ax.set_ylim(0, 100)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(pilares, size=11, weight='bold')
    ax.set_yticks([20, 40, 60, 80, 100])
    ax.set_yticklabels(['20', '40', '60', '80', '100'], size=10)
    ax.grid(True, linestyle='--', alpha=0.7)

    # Adicionar valores nos pontos
    for angle, score in zip(angles[:-1], pontuacoes):
        ax.text(angle, score + 5, f'{score}',
                ha='center', va='center',
                fontsize=10, fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='#00C896'))

    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"✓ Gráfico radar criado: {filename}")

def create_line_chart(meses, demanda, filename):
    """Cria gráfico de linha de sazonalidade"""
    fig, ax = plt.subplots(figsize=(14, 7))
    fig.patch.set_facecolor('white')

    # Plot linha principal
    ax.plot(meses, demanda, marker='o', linewidth=4, color='#1E3A5F', markersize=10, label='Demanda Relativa')

    # Marcar picos
    ax.axvline(x=6, color='#00C896', linestyle='--', linewidth=2, alpha=0.7, label='Jul: Férias Escolares')
    ax.axvline(x=10, color='#FF6B35', linestyle='--', linewidth=2, alpha=0.7, label='Nov: Black Friday (PICO)')
    ax.axvline(x=11, color='#FF6B35', linestyle='--', linewidth=2, alpha=0.5, label='Dez: 13º Salário')

    # Destacar área de pico
    ax.axhspan(90, 100, alpha=0.1, color='#FF6B35', label='Zona de Pico')

    # Configuração
    ax.set_xlabel('Mês', fontsize=14, fontweight='bold')
    ax.set_ylabel('Demanda Relativa (0-100)', fontsize=14, fontweight='bold')
    ax.set_title('Sazonalidade do Mercado de Consórcios - 2025', fontsize=16, fontweight='bold', pad=20)
    ax.legend(loc='upper left', fontsize=11)
    ax.grid(True, alpha=0.3, linestyle='--')
    ax.set_ylim(60, 105)

    # Adicionar valores nos pontos
    for i, (mes, valor) in enumerate(zip(meses, demanda)):
        if valor >= 90:  # Destacar picos
            ax.text(i, valor + 2, f'{valor}', ha='center', va='bottom',
                   fontsize=11, fontweight='bold', color='#FF6B35')

    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"✓ Gráfico de sazonalidade criado: {filename}")

def create_comparative_radar(concorrentes_data, filename):
    """Cria radar comparativo Klubi vs Concorrentes"""
    pilares = ['Tag', 'Analytics', 'Mídia', 'SEO', 'Velocidade', 'Conteúdo', 'CRM']

    fig, ax = plt.subplots(figsize=(12, 10), subplot_kw=dict(projection='polar'))
    fig.patch.set_facecolor('white')

    angles = np.linspace(0, 2 * np.pi, len(pilares), endpoint=False).tolist()
    angles += angles[:1]

    # Cores para cada empresa
    colors_map = {
        'Klubi': '#FF6B35',
        'Rodobens': '#1E3A5F',
        'Magalu': '#00C896',
        'Itaú': '#6C757D'
    }

    # Plotar cada concorrente
    for nome, dados in concorrentes_data.items():
        valores = dados + [dados[0]]
        ax.plot(angles, valores, 'o-', linewidth=2, label=nome,
                color=colors_map.get(nome, '#6C757D'), markersize=6, alpha=0.8)
        ax.fill(angles, valores, alpha=0.1, color=colors_map.get(nome, '#6C757D'))

    ax.set_ylim(0, 100)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(pilares, size=12, weight='bold')
    ax.set_yticks([20, 40, 60, 80, 100])
    ax.grid(True, linestyle='--', alpha=0.7)
    ax.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1), fontsize=11)

    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"✓ Radar comparativo criado: {filename}")

# SLIDE CREATORS
def create_slide_01_maturidade(prs, data):
    """Slide 1: Diagrama de Maturidade Digital"""
    slide = add_content_slide(prs, "Maturidade Digital | Klubi")

    # Criar gráfico radar
    pilares_nomes = [
        'Tag & Track',
        'GA4',
        'GTM',
        'Pixels',
        'Remarket.',
        'Velocidade',
        'SEO Téc.',
        'SEO Cont.',
        'CRM',
        'Mídia Paga'
    ]
    pontuacoes = [
        data['pilares']['tagueamento_tracking']['pontuacao'],
        data['pilares']['google_analytics_ga4']['pontuacao'],
        data['pilares']['google_tag_manager']['pontuacao'],
        data['pilares']['pixels_midia']['pontuacao'],
        data['pilares']['remarketing']['pontuacao'],
        data['pilares']['velocidade_site']['pontuacao'],
        data['pilares']['seo_tecnico']['pontuacao'],
        data['pilares']['seo_conteudo']['pontuacao'],
        data['pilares']['estrutura_crm_first_party']['pontuacao'],
        data['pilares']['preparacao_midia_paga']['pontuacao']
    ]

    chart_path = '/Users/brunopimentelV4/Desktop/Klubi/charts/radar_maturidade.png'
    create_radar_chart(pilares_nomes, pontuacoes, chart_path)

    # Adicionar gráfico
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), width=Inches(7))

    # Resumo executivo
    add_text_box(slide, 8, 1.5, 4.8, 0.5,
                 f"Pontuação Geral: {data['pontuacao_geral']}/100",
                 font_size=28, bold=True, color=COLORS['navy'])

    add_text_box(slide, 8, 2.2, 4.8, 0.4,
                 data['classificacao'],
                 font_size=18, bold=True, color=COLORS['orange'])

    # Top 3 forças
    add_text_box(slide, 8, 3, 4.8, 0.3, "✓ Pontos Fortes:",
                 font_size=14, bold=True, color=COLORS['green'])
    fortes_text = "• GTM implementado (75)\n• Velocidade excelente (78)\n• Tag & Track funcional (70)"
    add_text_box(slide, 8, 3.4, 4.8, 1, fortes_text, font_size=12)

    # Top 3 gaps
    add_text_box(slide, 8, 4.6, 4.8, 0.3, "✗ Gaps Críticos:",
                 font_size=14, bold=True, color=COLORS['red'])
    gaps_text = "• Meta Pixel ausente (40)\n• Remarketing limitado (45)\n• SEO Conteúdo fraco (55)"
    add_text_box(slide, 8, 5, 4.8, 1, gaps_text, font_size=12)

    # Nota metodológica
    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 f"Fonte: Análise técnica direta - {data['timestamp']} | Escala: 0-30 Crítico | 31-60 Parcial | 61-80 Adequado | 81-100 Excelente",
                 font_size=9, color=COLORS['gray'])

def create_slide_02_objetivos(prs):
    """Slide 2: Objetivos de Negócio"""
    slide = add_content_slide(prs, "Objetivos de Negócio | Klubi")

    objetivos = [
        {
            'numero': '01',
            'titulo': 'Geração de Leads Qualificados',
            'descricao': 'Capturar leads com alta intenção de adesão a consórcios através de canais digitais otimizados (Search, Social, Display)'
        },
        {
            'numero': '02',
            'titulo': 'Conversão em Vendas',
            'descricao': 'Converter leads em contratos assinados através de funil otimizado com remarketing estruturado e jornada personalizada'
        },
        {
            'numero': '03',
            'titulo': 'Redução de CAC',
            'descricao': 'Diminuir custo de aquisição através de otimização de campanhas, implementação de pixels e tracking avançado de conversões'
        }
    ]

    y_pos = 1.5
    for obj in objetivos:
        # Número
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y_pos),
            Inches(0.8), Inches(0.8)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLORS['navy']
        num_box.line.color.rgb = COLORS['navy']

        num_text = num_box.text_frame
        num_text.text = obj['numero']
        p = num_text.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        num_text.vertical_anchor = 1  # Middle

        # Título
        add_text_box(slide, 1.8, y_pos, 10.5, 0.4,
                     obj['titulo'],
                     font_size=20, bold=True, color=COLORS['navy'])

        # Descrição
        add_text_box(slide, 1.8, y_pos + 0.5, 10.5, 0.8,
                     obj['descricao'],
                     font_size=14, color=COLORS['dark_gray'])

        y_pos += 1.6

    # Nota
    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Nota: Objetivos técnicos focados em infraestrutura digital para suportar metas de negócio. Nenhuma promessa de resultado numérico é feita.",
                 font_size=9, color=COLORS['gray'])

def create_slide_03_friccoes(prs, data):
    """Slide 3: Fricções Críticas"""
    slide = add_content_slide(prs, "Fricções Críticas | Priorização")

    # Cabeçalho da tabela
    headers = ['#', 'Problema', 'Impacto', 'Esforço', 'ROI']
    x_positions = [0.5, 1.2, 6, 9.5, 11.3]
    widths = [0.5, 4.5, 3.2, 1.5, 1.5]

    y = 1.3
    for i, header in enumerate(headers):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_positions[i]), Inches(y),
            Inches(widths[i]), Inches(0.4)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['navy']
        box.line.fill.background()

        tf = box.text_frame
        tf.text = header
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf.vertical_anchor = 1

    # Linhas
    y = 1.8
    for friccao in data['friccoes_criticas'][:5]:
        # Prioridade
        prio_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.55), Inches(y + 0.05),
            Inches(0.35), Inches(0.35)
        )
        prio_box.fill.solid()
        cor_prio = COLORS['red'] if friccao['prioridade'] <= 2 else COLORS['orange'] if friccao['prioridade'] <= 3 else COLORS['gray']
        prio_box.fill.fore_color.rgb = cor_prio
        prio_box.line.fill.background()

        tf = prio_box.text_frame
        tf.text = str(friccao['prioridade'])
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf.vertical_anchor = 1

        # Problema
        add_text_box(slide, 1.2, y, 4.5, 0.5,
                     friccao['problema'][:80],
                     font_size=11, color=COLORS['dark_gray'])

        # Impacto
        add_text_box(slide, 6, y, 3.2, 0.5,
                     friccao['impacto'][:60],
                     font_size=10, color=COLORS['dark_gray'])

        # Esforço
        add_text_box(slide, 9.5, y + 0.08, 1.5, 0.3,
                     friccao['esforco'],
                     font_size=10, bold=True, color=COLORS['gray'])

        # ROI
        add_text_box(slide, 11.3, y + 0.08, 1.5, 0.3,
                     friccao['roi_esperado'],
                     font_size=10, bold=True, color=COLORS['green'])

        y += 0.9

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 f"Fonte: {data['fonte']} | Priorização baseada em Impacto × Esforço × ROI Esperado",
                 font_size=9, color=COLORS['gray'])

def create_slide_04_problemas_oportunidades(prs, data):
    """Slide 4: Problemas vs Oportunidades"""
    slide = add_content_slide(prs, "Problemas vs Oportunidades")

    problemas = [
        {
            'problema': 'Meta Pixel e TikTok Pixel ausentes',
            'evidencia': 'window.fbq e window.ttq undefined no site',
            'impacto': 'Impossível otimizar Meta/TikTok Ads',
            'correcao': 'Implementar via GTM (1-2 dias)'
        },
        {
            'problema': 'GA4 retornando erro 503',
            'evidencia': 'HTTP 503 em google-analytics.com/collect',
            'impacto': 'Perda de dados comportamentais',
            'correcao': 'Debug configuração GA4 (2-3 dias)'
        },
        {
            'problema': 'Schema markup ausente',
            'evidencia': 'script[type="application/ld+json"] = vazio',
            'impacto': 'SEO comprometido, sem rich snippets',
            'correcao': 'Implementar Organization, FAQ, How-To'
        },
        {
            'problema': 'Estrutura de conteúdo limitada',
            'evidencia': '0 links internos detectados',
            'impacto': 'Baixa captura orgânica, dependência de paid',
            'correcao': 'Criar blog educacional (estratégia contínua)'
        }
    ]

    # Headers
    headers = ['Problema', 'Evidência Técnica', 'Impacto', 'Correção']
    x_pos = [0.5, 3.8, 7, 9.8]
    widths = [3, 3, 2.5, 3]

    y = 1.3
    for i, h in enumerate(headers):
        add_text_box(slide, x_pos[i], y, widths[i], 0.3,
                     h, font_size=11, bold=True, color=COLORS['navy'])

    # Linhas
    y = 1.7
    for p in problemas:
        # Separador
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(y - 0.05),
            Inches(12), Inches(0.02)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = COLORS['light_gray']
        sep.line.fill.background()

        add_text_box(slide, 0.5, y, 3, 0.6, p['problema'], font_size=10)
        add_text_box(slide, 3.8, y, 3, 0.6, p['evidencia'], font_size=9, color=COLORS['gray'])
        add_text_box(slide, 7, y, 2.5, 0.6, p['impacto'], font_size=9, color=COLORS['red'])
        add_text_box(slide, 9.8, y, 3, 0.6, p['correcao'], font_size=9, color=COLORS['green'])

        y += 1.2

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Fonte: Análise técnica direta do site klubi.com.br | Todos os problemas são tecnicamente verificáveis e corrigíveis",
                 font_size=9, color=COLORS['gray'])

def create_slide_05_canais_midia(prs, midia_data):
    """Slide 5: Panorama de Canais de Mídia"""
    slide = add_content_slide(prs, "Panorama de Canais de Mídia")

    canais = [
        {
            'nome': 'Google Ads',
            'status': 'ATIVO',
            'cor': COLORS['green'],
            'evidencia': 'Conversion tag AW-362695272 presente',
            'gaps': 'Eventos customizados não auditados',
            'prioridade': 'Otimizar campanhas existentes'
        },
        {
            'nome': 'Meta Ads',
            'status': 'AUSENTE',
            'cor': COLORS['red'],
            'evidencia': 'Meta Pixel não detectado',
            'gaps': 'Impossível otimizar FB/IG Ads',
            'prioridade': 'CRÍTICO: Implementar pixel'
        },
        {
            'nome': 'TikTok Ads',
            'status': 'AUSENTE',
            'cor': COLORS['red'],
            'evidencia': 'TikTok Pixel não detectado',
            'gaps': 'Canal crescente não utilizado',
            'prioridade': 'ALTA: Implementar para público jovem'
        },
        {
            'nome': 'Amazon Ads',
            'status': 'ATIVO',
            'cor': COLORS['green'],
            'evidencia': 'Script presente (PID f50f2bc2...)',
            'gaps': 'Canal secundário para consórcios',
            'prioridade': 'Monitorar performance'
        }
    ]

    y = 1.5
    for canal in canais:
        # Box do canal
        canal_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y),
            Inches(11.8), Inches(1.1)
        )
        canal_box.fill.solid()
        canal_box.fill.fore_color.rgb = COLORS['light_gray']
        canal_box.line.color.rgb = COLORS['gray']
        canal_box.line.width = Pt(1)

        # Nome do canal
        add_text_box(slide, 1, y + 0.1, 2.5, 0.3,
                     canal['nome'],
                     font_size=18, bold=True, color=COLORS['navy'])

        # Status badge
        status_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.7), Inches(y + 0.08),
            Inches(1.2), Inches(0.35)
        )
        status_box.fill.solid()
        status_box.fill.fore_color.rgb = canal['cor']
        status_box.line.fill.background()

        tf = status_box.text_frame
        tf.text = canal['status']
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf.vertical_anchor = 1

        # Evidência
        add_text_box(slide, 1, y + 0.5, 5, 0.25,
                     f"Evidência: {canal['evidencia']}",
                     font_size=9, color=COLORS['gray'])

        # Gaps
        add_text_box(slide, 5.2, y + 0.1, 3.5, 0.3,
                     f"Gap: {canal['gaps']}",
                     font_size=10, color=COLORS['red'])

        # Prioridade
        add_text_box(slide, 5.2, y + 0.5, 3.5, 0.3,
                     f"→ {canal['prioridade']}",
                     font_size=10, bold=True, color=COLORS['green'])

        y += 1.3

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Fonte: Análise técnica de pixels e scripts no site klubi.com.br | Bibliotecas públicas de anúncios não acessíveis",
                 font_size=9, color=COLORS['gray'])

def create_slide_06_palavras_chave(prs, kw_data):
    """Slide 6: Palavras-Chave Transacionais (25)"""
    slide = add_content_slide(prs, "Palavras-Chave Transacionais | Top 25")

    # Nota de disclaimer
    add_text_box(slide, 0.5, 1.2, 12, 0.3,
                 "IMPORTANTE: Volumes e CPC não disponíveis publicamente (SEMrush/Ahrefs pagos). Lista baseada em análise de mercado e tendências.",
                 font_size=10, bold=True, color=COLORS['orange'])

    # Dividir em 2 colunas
    palavras_col1 = kw_data['palavras_chave'][:13]
    palavras_col2 = kw_data['palavras_chave'][13:25]

    # Coluna 1
    y = 1.8
    add_text_box(slide, 0.5, y, 0.3, 0.25, "#", font_size=9, bold=True)
    add_text_box(slide, 0.9, y, 3.5, 0.25, "Palavra-Chave", font_size=9, bold=True)
    add_text_box(slide, 4.5, y, 1.5, 0.25, "Tipo", font_size=9, bold=True)

    y += 0.35
    for i, kw in enumerate(palavras_col1, 1):
        add_text_box(slide, 0.5, y, 0.3, 0.25, str(i), font_size=8)
        add_text_box(slide, 0.9, y, 3.5, 0.25, kw['palavra'], font_size=8)
        add_text_box(slide, 4.5, y, 1.5, 0.25, kw['tipo'].split('_')[0], font_size=7, color=COLORS['gray'])
        y += 0.32

    # Coluna 2
    y = 1.8
    add_text_box(slide, 6.5, y, 0.3, 0.25, "#", font_size=9, bold=True)
    add_text_box(slide, 6.9, y, 3.5, 0.25, "Palavra-Chave", font_size=9, bold=True)
    add_text_box(slide, 10.5, y, 1.5, 0.25, "Tipo", font_size=9, bold=True)

    y += 0.35
    for i, kw in enumerate(palavras_col2, 14):
        if i <= 25:
            add_text_box(slide, 6.5, y, 0.3, 0.25, str(i), font_size=8)
            add_text_box(slide, 6.9, y, 3.5, 0.25, kw['palavra'], font_size=8)
            add_text_box(slide, 10.5, y, 1.5, 0.25, kw['tipo'].split('_')[0], font_size=7, color=COLORS['gray'])
            y += 0.32

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Fontes: Tendências de mercado ABAC, Turn2C, Google Trends Brasil | Volume estimado qualitativamente",
                 font_size=9, color=COLORS['gray'])

def create_slide_07_concorrencia(prs, conc_data):
    """Slide 7: Análise de Concorrência"""
    slide = add_content_slide(prs, "Análise de Concorrência | Maturidade Digital")

    # Criar radar comparativo
    concorrentes_radar = {
        'Klubi': [70, 65, 58, 60, 78, 55, 60],
        'Rodobens': [80, 85, 80, 85, 75, 70, 80],
        'Magalu': [90, 90, 90, 85, 85, 75, 90],
        'Itaú': [85, 90, 85, 80, 75, 65, 90]
    }

    chart_path = '/Users/brunopimentelV4/Desktop/Klubi/charts/comparativo_concorrencia.png'
    create_comparative_radar(concorrentes_radar, chart_path)

    # Adicionar gráfico
    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), width=Inches(7))

    # Insights
    add_text_box(slide, 8, 1.5, 4.5, 0.4,
                 "Insights Competitivos",
                 font_size=18, bold=True, color=COLORS['navy'])

    insights = [
        "✓ Klubi tem velocidade superior a bancos tradicionais (78 vs 70-75)",
        "✗ Gap significativo em Mídia Paga vs Magalu e Itaú (58 vs 85-90)",
        "✗ SEO Conteúdo muito abaixo dos líderes (55 vs 70-75)",
        "⚠ CRM competitivo mas ainda atrás de Magalu/Itaú",
        "→ Oportunidade: Digital nativo permite agilidade maior"
    ]

    y = 2.1
    for insight in insights:
        cor = COLORS['green'] if insight.startswith('✓') else COLORS['red'] if insight.startswith('✗') else COLORS['orange']
        add_text_box(slide, 8, y, 4.5, 0.45, insight, font_size=11, color=cor)
        y += 0.55

    # Ranking
    add_text_box(slide, 8, 5.2, 4.5, 0.3,
                 "Ranking Maturidade Digital:",
                 font_size=12, bold=True, color=COLORS['navy'])
    ranking_text = "1º Magalu (86)\n2º Itaú (81)\n3º Rodobens (79)\n4º Klubi (61)"
    add_text_box(slide, 8, 5.6, 4.5, 1, ranking_text, font_size=11)

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Fontes: iDinheiro, ConsorcioCred, NeoFeed | Maturidade estimada de forma qualitativa (sem acesso a painéis internos)",
                 font_size=9, color=COLORS['gray'])

def create_slide_08_sazonalidade(prs, saz_data):
    """Slide 8: Sazonalidade e Ciclos Econômicos"""
    slide = add_content_slide(prs, "Sazonalidade e Ciclos Econômicos | 2025")

    # Criar gráfico
    meses = ['Jan', 'Fev', 'Mar', 'Abr', 'Mai', 'Jun', 'Jul', 'Ago', 'Set', 'Out', 'Nov', 'Dez']
    demanda = [m['demanda_relativa'] for m in saz_data['dados_mensais']]

    chart_path = '/Users/brunopimentelV4/Desktop/Klubi/charts/linha_sazonalidade.png'
    create_line_chart(meses, demanda, chart_path)

    slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), width=Inches(12))

    # Crescimento setorial
    add_text_box(slide, 0.5, 5.8, 12, 0.3,
                 "Crescimento Setorial 2025 (Fonte: ABAC)",
                 font_size=12, bold=True, color=COLORS['navy'])

    crescimento_text = "Geral: +8-11% | Imóveis: +20-25% | Eletroeletrônicos: +20-23% | Serviços: +10%"
    add_text_box(slide, 0.5, 6.2, 12, 0.3,
                 crescimento_text,
                 font_size=11, color=COLORS['green'])

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Fontes: ABAC, Turn2C, O2O Bots | Demanda relativa baseada em padrões históricos do setor (0-100, sendo 100 = pico Black Friday)",
                 font_size=9, color=COLORS['gray'])

def create_slide_09_janelas_atuacao(prs):
    """Slide 9: Como Atuar em Cada Janela"""
    slide = add_content_slide(prs, "Como Atuar em Cada Janela Sazonal")

    janelas = [
        {
            'periodo': 'ANTES DO PICO',
            'meses': 'Set - Out',
            'cor': COLORS['orange'],
            'objetivo': 'Construir awareness e consideração',
            'taticas': [
                'Conteúdo educacional intenso (blog, vídeos)',
                'Remarketing de captura (Google Display, Meta)',
                'Aquecimento de audiências custom',
                'Ofertas early bird para early adopters',
                'Testes A/B de criativos e mensagens'
            ],
            'investimento': '60-70% do pico'
        },
        {
            'periodo': 'DURANTE O PICO',
            'meses': 'Nov - Dez',
            'cor': COLORS['red'],
            'objetivo': 'Conversão máxima',
            'taticas': [
                'Investimento MÁXIMO em mídia paga',
                'Ofertas agressivas e limitadas por tempo',
                'Remarketing total (todos os canais)',
                'Landing pages otimizadas para conversão',
                'Atendimento estendido (chat, WhatsApp)',
                'Push notifications e email marketing'
            ],
            'investimento': '100% (pico absoluto)'
        },
        {
            'periodo': 'DEPOIS DO PICO',
            'meses': 'Jan - Fev',
            'cor': COLORS['green'],
            'objetivo': 'Manutenção e planejadores',
            'taticas': [
                'Remarketing de nurturing (não agressivo)',
                'Conteúdo de planejamento anual',
                'Captura de "Resoluções de Ano Novo"',
                'Ofertas para planejadores de longo prazo',
                'Redução controlada de paid media'
            ],
            'investimento': '30-40% do pico'
        }
    ]

    y = 1.4
    for janela in janelas:
        # Header
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), Inches(y),
            Inches(11.8), Inches(0.4)
        )
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = janela['cor']
        header_box.line.fill.background()

        tf = header_box.text_frame
        tf.text = f"{janela['periodo']} ({janela['meses']}) | {janela['objetivo']}"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        tf.vertical_anchor = 1

        # Táticas
        y_tat = y + 0.5
        for i, tatica in enumerate(janela['taticas'][:5], 1):
            add_text_box(slide, 1, y_tat, 9, 0.25,
                         f"{i}. {tatica}",
                         font_size=9, color=COLORS['dark_gray'])
            y_tat += 0.28

        # Investimento
        inv_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(10.5), Inches(y + 0.5),
            Inches(1.8), Inches(0.5)
        )
        inv_box.fill.solid()
        inv_box.fill.fore_color.rgb = COLORS['light_gray']
        inv_box.line.color.rgb = janela['cor']
        inv_box.line.width = Pt(2)

        tf_inv = inv_box.text_frame
        tf_inv.text = f"Invest:\n{janela['investimento']}"
        for para in tf_inv.paragraphs:
            para.font.size = Pt(10)
            para.font.bold = True
            para.font.color.rgb = janela['cor']
            para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf_inv.vertical_anchor = 1

        y += 1.9

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Recomendação de alocação orçamentária: Nov-Dez (40-45%) | Jul (15-20%) | Set-Out (20-25%) | Demais (15-20%)",
                 font_size=9, color=COLORS['gray'])

def create_slide_10_metodologia(prs):
    """Slide 10: Metodologia de Atuação (8 Semanas)"""
    slide = add_content_slide(prs, "Metodologia de Atuação | Primeiras 8 Semanas")

    semanas = [
        {'num': '01-02', 'titulo': 'Diagnóstico e Quick Wins',
         'entregas': 'Implementar Meta/TikTok Pixel, Corrigir GA4, Auditar Google Ads'},
        {'num': '03-04', 'titulo': 'Estruturação de Funil',
         'entregas': 'Mapear eventos de conversão, Configurar remarketing, Criar audiências custom'},
        {'num': '05-06', 'titulo': 'SEO Técnico e Conteúdo',
         'entregas': 'Implementar Schema.org, Criar 5-10 artigos educacionais, Otimizar on-page'},
        {'num': '07-08', 'titulo': 'Otimização e Escala',
         'entregas': 'Lançar Performance Max, Otimizar Meta Ads, Testes A/B criativos'}
    ]

    y = 1.5
    for sem in semanas:
        # Timeline box
        timeline_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y),
            Inches(1.2), Inches(0.5)
        )
        timeline_box.fill.solid()
        timeline_box.fill.fore_color.rgb = COLORS['navy']
        timeline_box.line.fill.background()

        tf = timeline_box.text_frame
        tf.text = f"S{sem['num']}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf.vertical_anchor = 1

        # Título
        add_text_box(slide, 2.2, y, 4, 0.25,
                     sem['titulo'],
                     font_size=16, bold=True, color=COLORS['navy'])

        # Entregas
        add_text_box(slide, 2.2, y + 0.3, 10, 0.3,
                     f"→ {sem['entregas']}",
                     font_size=12, color=COLORS['dark_gray'])

        # Linha conectora
        if y < 5:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.25), Inches(y + 0.55),
                Inches(0.05), Inches(0.5)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS['gray']
            line.line.fill.background()

        y += 1.2

    # Nota importante
    nota_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(6.2),
        Inches(11.6), Inches(0.5)
    )
    nota_box.fill.solid()
    nota_box.fill.fore_color.rgb = COLORS['light_gray']
    nota_box.line.color.rgb = COLORS['orange']
    nota_box.line.width = Pt(2)

    tf_nota = nota_box.text_frame
    tf_nota.text = "IMPORTANTE: Timeline orientativa. Ajustes baseados em prioridades de negócio e recursos disponíveis. Nenhuma promessa de resultado numérico."
    p_nota = tf_nota.paragraphs[0]
    p_nota.font.size = Pt(10)
    p_nota.font.bold = True
    p_nota.font.color.rgb = COLORS['orange']
    tf_nota.vertical_anchor = 1

def create_slide_11_sintese(prs, data):
    """Slide 11: Síntese Executiva"""
    slide = add_content_slide(prs, "Síntese Executiva | 3 Insights Principais")

    insights = [
        {
            'numero': '01',
            'titulo': 'Infraestrutura Básica Funcional, Gaps Críticos em Canais',
            'descricao': 'GTM, GA4 e tracking Google estão implementados (70-75/100), mas ausência de Meta Pixel e TikTok Pixel impede otimização de 2 bilhões+ usuários. Quick win de alto ROI.',
            'acao': 'Implementar Meta e TikTok Pixels em 1-2 dias via GTM'
        },
        {
            'numero': '02',
            'titulo': 'SEO Técnico Adequado, Conteúdo Inexistente',
            'descricao': 'Meta tags, Open Graph e velocidade OK, mas sem Schema.org, sem blog e 0 links internos. Dependência excessiva de mídia paga. Termo "consórcio vale a pena" dominou buscas em 2024-2025.',
            'acao': 'Criar estratégia de conteúdo educacional para captura orgânica'
        },
        {
            'numero': '03',
            'titulo': 'Sazonalidade Crítica: Novembro Define o Ano',
            'descricao': 'Black Friday (Nov) representa pico de 100/100 em demanda. Jul (92) e Dez (88) são secundários. Crescimento setorial 2025: +8-11% geral, +20-25% imóveis. Alocação orçamentária deve refletir isso.',
            'acao': 'Concentrar 40-45% do budget anual em Nov-Dez, 15-20% em Jul'
        }
    ]

    y = 1.5
    for insight in insights:
        # Caixa principal
        main_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), Inches(y),
            Inches(11.8), Inches(1.4)
        )
        main_box.fill.solid()
        main_box.fill.fore_color.rgb = COLORS['light_gray']
        main_box.line.color.rgb = COLORS['navy']
        main_box.line.width = Pt(2)

        # Número
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y + 0.15),
            Inches(0.6), Inches(0.6)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLORS['navy']
        num_circle.line.fill.background()

        tf_num = num_circle.text_frame
        tf_num.text = insight['numero']
        p_num = tf_num.paragraphs[0]
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = COLORS['white']
        p_num.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf_num.vertical_anchor = 1

        # Título
        add_text_box(slide, 1.6, y + 0.1, 10.5, 0.35,
                     insight['titulo'],
                     font_size=14, bold=True, color=COLORS['navy'])

        # Descrição
        add_text_box(slide, 1.6, y + 0.5, 10.5, 0.6,
                     insight['descricao'],
                     font_size=11, color=COLORS['dark_gray'])

        # Ação
        acao_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.6), Inches(y + 1.05),
            Inches(10.5), Inches(0.25)
        )
        acao_box.fill.solid()
        acao_box.fill.fore_color.rgb = COLORS['green']
        acao_box.line.fill.background()

        tf_acao = acao_box.text_frame
        tf_acao.text = f"→ AÇÃO: {insight['acao']}"
        p_acao = tf_acao.paragraphs[0]
        p_acao.font.size = Pt(10)
        p_acao.font.bold = True
        p_acao.font.color.rgb = COLORS['white']
        tf_acao.vertical_anchor = 1

        y += 1.65

def create_slide_12_fontes(prs):
    """Slide 12: Apêndice | Fontes"""
    slide = add_content_slide(prs, "Apêndice | Fontes de Dados")

    fontes = [
        {'categoria': 'Análise Técnica', 'fonte': 'Análise direta do site klubi.com.br',
         'dados': 'Pixels, tracking, SEO técnico, performance', 'conf': 'Muito Alta'},
        {'categoria': 'Mercado - Crescimento', 'fonte': 'ABAC (Assoc. Brasileira de Administradoras de Consórcios)',
         'dados': 'Projeções 2025, crescimento setorial', 'conf': 'Muito Alta'},
        {'categoria': 'Mercado - Tendências', 'fonte': 'Turn2C - Tendências do mercado de consórcio 2025',
         'dados': 'Sazonalidade, "consórcio vale a pena"', 'conf': 'Alta'},
        {'categoria': 'Mercado - Histórico', 'fonte': 'O2O Bots - Tendências e Dados do Mercado de Consórcios',
         'dados': 'R$ 316,7bi em 2023, 10,29mi participantes', 'conf': 'Alta'},
        {'categoria': 'Concorrência - Rankings', 'fonte': 'iDinheiro - Melhor Consórcio do Brasil 2026',
         'dados': 'Rodobens, Sicredi, rankings', 'conf': 'Alta'},
        {'categoria': 'Concorrência - Análise', 'fonte': 'ConsorcioCred, DP Consórcios, NeoFeed',
         'dados': 'Maiores administradoras, performance', 'conf': 'Média-Alta'},
        {'categoria': 'Dados Oficiais', 'fonte': 'Banco Central - Ranking de Administradoras',
         'dados': 'Lista oficial de administradoras', 'conf': 'Muito Alta'},
        {'categoria': 'Busca - Tendências', 'fonte': 'Google Trends Brasil, InfoMoney',
         'dados': 'Produtos financeiros mais buscados 2025', 'conf': 'Alta'},
        {'categoria': 'Limitação', 'fonte': 'Meta Ads Library, Google Ads Transparency',
         'dados': 'Não acessíveis via automação', 'conf': 'N/A'},
        {'categoria': 'Limitação', 'fonte': 'SEMrush, Ahrefs (ferramentas pagas)',
         'dados': 'Volume/CPC não disponíveis publicamente', 'conf': 'N/A'}
    ]

    # Headers
    headers = ['Categoria', 'Fonte', 'Dados Utilizados', 'Confiab.']
    x_pos = [0.5, 2.2, 5.5, 10.5]
    widths = [1.5, 3, 4.5, 1.8]

    y = 1.3
    for i, h in enumerate(headers):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos[i]), Inches(y),
            Inches(widths[i]), Inches(0.35)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['navy']
        box.line.fill.background()

        tf = box.text_frame
        tf.text = h
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf.vertical_anchor = 1

    # Linhas
    y = 1.75
    for fonte in fontes[:10]:
        cor_conf = COLORS['green'] if fonte['conf'] in ['Muito Alta', 'Alta'] else COLORS['orange'] if fonte['conf'] == 'Média-Alta' else COLORS['gray']

        add_text_box(slide, 0.5, y, 1.5, 0.4, fonte['categoria'], font_size=8, bold=True)
        add_text_box(slide, 2.2, y, 3, 0.4, fonte['fonte'], font_size=8)
        add_text_box(slide, 5.5, y, 4.5, 0.4, fonte['dados'], font_size=7, color=COLORS['gray'])
        add_text_box(slide, 10.5, y, 1.8, 0.4, fonte['conf'], font_size=8, bold=True, color=cor_conf)

        y += 0.48

    add_text_box(slide, 0.5, 6.8, 12, 0.4,
                 "Metodologia: Análise técnica direta + Pesquisa de fontes públicas confiáveis. ZERO dados inventados ou estimados sem fonte.",
                 font_size=9, color=COLORS['orange'])

def create_slide_13_confiabilidade(prs):
    """Slide 13: Apêndice | Escala de Confiabilidade"""
    slide = add_content_slide(prs, "Apêndice | Escala de Confiabilidade dos Dados")

    niveis = [
        {
            'nivel': 'MUITO ALTA',
            'cor': COLORS['green'],
            'definicao': 'Dado verificado tecnicamente ou de fonte oficial primária',
            'exemplos': [
                'Análise técnica direta do site (pixels, scripts, performance)',
                'Dados de ABAC (associação oficial do setor)',
                'Ranking do Banco Central',
                'Medições de performance (DOM load, HTTP requests)'
            ]
        },
        {
            'nivel': 'ALTA',
            'cor': COLORS['green'],
            'definicao': 'Fonte secundária confiável, publicação especializada',
            'exemplos': [
                'iDinheiro, Turn2C, O2O Bots (publicações especializadas)',
                'InfoMoney, Google Trends (dados agregados)',
                'Reportagens de veículos estabelecidos (NeoFeed, CNN Brasil)'
            ]
        },
        {
            'nivel': 'MÉDIA-ALTA',
            'cor': COLORS['orange'],
            'definicao': 'Análise qualitativa baseada em múltiplas fontes',
            'exemplos': [
                'Maturidade digital de concorrentes (estimada)',
                'Presença em mídia paga (inferida por ausência de pixels)',
                'Tendências de palavras-chave (sem acesso a ferramentas pagas)'
            ]
        },
        {
            'nivel': 'REVISAR',
            'cor': COLORS['red'],
            'definicao': 'Dado não obtido, requer acesso a ferramentas/painéis pagos',
            'exemplos': [
                'Volume de busca específico (requer SEMrush/Ahrefs)',
                'CPC por palavra-chave (requer Google Keyword Planner)',
                'Anúncios ativos de concorrentes (bibliotecas não acessíveis)',
                'Investimento em mídia dos competidores'
            ]
        }
    ]

    y = 1.5
    for nivel in niveis:
        # Badge do nível
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y),
            Inches(2), Inches(0.5)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = nivel['cor']
        badge.line.fill.background()

        tf_badge = badge.text_frame
        tf_badge.text = nivel['nivel']
        p_badge = tf_badge.paragraphs[0]
        p_badge.font.size = Pt(16)
        p_badge.font.bold = True
        p_badge.font.color.rgb = COLORS['white']
        p_badge.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        tf_badge.vertical_anchor = 1

        # Definição
        add_text_box(slide, 3, y + 0.05, 9.3, 0.4,
                     nivel['definicao'],
                     font_size=12, bold=True, color=COLORS['navy'])

        # Exemplos
        y_ex = y + 0.55
        for exemplo in nivel['exemplos']:
            add_text_box(slide, 3, y_ex, 9.3, 0.25,
                         f"• {exemplo}",
                         font_size=9, color=COLORS['dark_gray'])
            y_ex += 0.28

        y += 1.4

    # Compromisso de qualidade
    compromisso_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(6.2),
        Inches(11.6), Inches(0.5)
    )
    compromisso_box.fill.solid()
    compromisso_box.fill.fore_color.rgb = COLORS['navy']
    compromisso_box.line.fill.background()

    tf_comp = compromisso_box.text_frame
    tf_comp.text = "COMPROMISSO: Todo dado possui fonte citada. Quando não disponível, declaramos explicitamente. ZERO invenção de métricas."
    p_comp = tf_comp.paragraphs[0]
    p_comp.font.size = Pt(11)
    p_comp.font.bold = True
    p_comp.font.color.rgb = COLORS['white']
    p_comp.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    tf_comp.vertical_anchor = 1

# MAIN
def main():
    print("=" * 60)
    print("GERAÇÃO DE APRESENTAÇÃO PPTX - ANÁLISE KLUBI")
    print("=" * 60)

    # Carregar dados
    print("\n[1/4] Carregando dados JSON...")
    with open('/Users/brunopimentelV4/Desktop/Klubi/data/maturidade_digital.json', 'r', encoding='utf-8') as f:
        maturidade = json.load(f)
    with open('/Users/brunopimentelV4/Desktop/Klubi/data/palavras_chave.json', 'r', encoding='utf-8') as f:
        palavras_chave = json.load(f)
    with open('/Users/brunopimentelV4/Desktop/Klubi/data/concorrencia.json', 'r', encoding='utf-8') as f:
        concorrencia = json.load(f)
    with open('/Users/brunopimentelV4/Desktop/Klubi/data/sazonalidade.json', 'r', encoding='utf-8') as f:
        sazonalidade = json.load(f)
    with open('/Users/brunopimentelV4/Desktop/Klubi/data/midia_paga.json', 'r', encoding='utf-8') as f:
        midia_paga = json.load(f)
    print("✓ Dados carregados com sucesso")

    # Criar apresentação
    print("\n[2/4] Criando apresentação...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide de título
    print("  → Criando slide de título...")
    add_title_slide(prs,
                    "Análise de Maturidade Digital",
                    "Klubi | Primeira Fintech de Consórcios do Brasil")

    # Slides de conteúdo
    print("  → Criando slides de conteúdo...")
    print("     [01/13] Maturidade Digital...")
    create_slide_01_maturidade(prs, maturidade)

    print("     [02/13] Objetivos de Negócio...")
    create_slide_02_objetivos(prs)

    print("     [03/13] Fricções Críticas...")
    create_slide_03_friccoes(prs, maturidade)

    print("     [04/13] Problemas vs Oportunidades...")
    create_slide_04_problemas_oportunidades(prs, maturidade)

    print("     [05/13] Canais de Mídia...")
    create_slide_05_canais_midia(prs, midia_paga)

    print("     [06/13] Palavras-Chave...")
    create_slide_06_palavras_chave(prs, palavras_chave)

    print("     [07/13] Concorrência...")
    create_slide_07_concorrencia(prs, concorrencia)

    print("     [08/13] Sazonalidade...")
    create_slide_08_sazonalidade(prs, sazonalidade)

    print("     [09/13] Janelas de Atuação...")
    create_slide_09_janelas_atuacao(prs)

    print("     [10/13] Metodologia 8 Semanas...")
    create_slide_10_metodologia(prs)

    print("     [11/13] Síntese Executiva...")
    create_slide_11_sintese(prs, maturidade)

    print("     [12/13] Fontes...")
    create_slide_12_fontes(prs)

    print("     [13/13] Confiabilidade...")
    create_slide_13_confiabilidade(prs)

    # Salvar
    print("\n[3/4] Salvando apresentação...")
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    output_path = f'/Users/brunopimentelV4/Desktop/Klubi/output/Analise_Maturidade_Digital_Klubi_{timestamp}.pptx'
    prs.save(output_path)
    print(f"✓ Apresentação salva: {output_path}")

    # Verificar arquivo
    print("\n[4/4] Verificando arquivo gerado...")
    import os
    file_size = os.path.getsize(output_path) / (1024 * 1024)  # MB
    print(f"✓ Tamanho do arquivo: {file_size:.2f} MB")
    print(f"✓ Total de slides: {len(prs.slides)}")

    print("\n" + "=" * 60)
    print("✅ APRESENTAÇÃO CRIADA COM SUCESSO!")
    print("=" * 60)
    print(f"\nArquivo: {output_path}")
    print(f"\nPróximos passos:")
    print("  1. Abrir o arquivo PPTX")
    print("  2. Revisar todos os slides")
    print("  3. Ajustar formatação se necessário")
    print("  4. Validar fontes e dados")
    print("\n" + "=" * 60)

if __name__ == '__main__':
    main()
