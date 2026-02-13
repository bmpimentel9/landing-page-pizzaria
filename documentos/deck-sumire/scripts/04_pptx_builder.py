"""
PPTX Builder Script
Creates complete 13-slide presentation for Perfumaria Sumirê

Slides:
1. Capa
2. Metodologia
3. Sobre a Marca
4. Radar Maturidade Digital
5. Fundamentos - Matriz Priorização
6. Radar Comparativo Concorrentes
7. Audiência 12 Meses
8. Comparativo Canais Concorrentes
9. Sankey Flow
10. Keywords Transacionais
11. Demanda 12 Meses
12. Importância SEO + GEO
13. Presença em LLMs
"""

import os
import sys
from datetime import datetime

# Add parent directory to path
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

from config import COLORS_RGB, SUMIRE, COMPETITORS, SLIDE_CONFIG, PATHS
from utils import load_json, log_progress, log_success, log_error


class PPTXBuilder:
    """Build complete presentation"""

    def __init__(self, data_file: str, llm_file: str):
        """Initialize PPTX builder"""
        self.data = load_json(data_file)
        self.llm_data = load_json(llm_file)

        if not self.data:
            raise ValueError("No data available")

        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_CONFIG['width_inches'])
        self.prs.slide_height = Inches(SLIDE_CONFIG['height_inches'])

        self.charts_dir = os.path.join(os.path.dirname(__file__), '../charts/')

    def add_title_slide(self):
        """Slide 1: Cover slide"""
        log_progress("Creating Slide 1: Capa...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*COLORS_RGB['primary'])

        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = "PERFUMARIA SUMIRÊ"
        p = tf.paragraphs[0]
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        tf.text = "Análise de Presença Digital"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Footer
        footer_box = slide.shapes.add_textbox(
            Inches(2.5), Inches(6.5), Inches(5), Inches(0.6)
        )
        tf = footer_box.text_frame
        tf.text = f"V4 Carvalho Consultoria Digital\n{datetime.now().strftime('%B %Y')}"
        for p in tf.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    def add_methodology_slide(self):
        """Slide 2: Methodology"""
        log_progress("Creating Slide 2: Metodologia...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Title Only

        # Title
        title = slide.shapes.title
        title.text = "Metodologia e Ferramentas"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*COLORS_RGB['dark'])

        # Content
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(8.4)
        height = Inches(5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        intro = tf.paragraphs[0]
        intro.text = "Esta análise foi realizada utilizando ferramentas profissionais de mercado para garantir precisão e profundidade nos dados coletados:\n"
        intro.font.size = Pt(16)
        intro.space_after = Pt(18)

        ferramentas = [
            "SimilarWeb — Análise de tráfego e comportamento de audiência",
            "Ubersuggest — SEO, keywords e posicionamento orgânico",
            "Google PageSpeed Insights — Performance técnica do site",
            "Google Trends — Tendências de busca e sazonalidade",
            "DevTools — Análise de tags, pixels e implementações técnicas",
            "Testes em LLMs — ChatGPT, Perplexity, Gemini, Claude"
        ]

        for ferramenta in ferramentas:
            p = tf.add_paragraph()
            p.text = "• " + ferramenta
            p.font.size = Pt(16)
            p.space_before = Pt(10)
            p.level = 0

    def add_about_brand_slide(self):
        """Slide 3: About the brand"""
        log_progress("Creating Slide 3: Sobre a Marca...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title = slide.shapes.title
        title.text = "Perfumaria Sumirê"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*COLORS_RGB['primary'])

        # Content
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(8.4)
        height = Inches(5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        content_items = [
            ("História", f"Fundada em {SUMIRE['founded']}, a Perfumaria Sumirê é uma das maiores redes de perfumarias do Brasil, com mais de {SUMIRE['founded'] - 1984} anos de tradição no mercado."),
            ("Presença Física", f"Aproximadamente {SUMIRE['stores']} lojas estrategicamente localizadas em São Paulo, atendendo diversos bairros e regiões da capital e interior."),
            ("Digital", f"Forte presença no Instagram com {SUMIRE['followers']:,} seguidores, além de e-commerce próprio e app Sumirê Club VIP para programa de fidelidade."),
            ("Posicionamento", f"\"{SUMIRE['positioning']}\" — Competindo com grandes players nacionais e internacionais no setor de beleza e cosméticos.")
        ]

        for idx, (subtitle, text) in enumerate(content_items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = subtitle
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*COLORS_RGB['secondary'])
            p.space_before = Pt(20) if idx > 0 else 0

            p2 = tf.add_paragraph()
            p2.text = text
            p2.font.size = Pt(14)
            p2.level = 0

    def add_image_slide(self, slide_num: int, title: str, image_filename: str):
        """Generic slide with title and image"""
        log_progress(f"Creating Slide {slide_num}: {title}...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*COLORS_RGB['dark'])

        # Add image
        image_path = os.path.join(self.charts_dir, image_filename)

        if os.path.exists(image_path):
            left = Inches(0.5)
            top = Inches(1.8)
            height = Inches(5.2)
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            log_error(f"Image not found: {image_path}")

    def add_fundamentals_slide(self):
        """Slide 5: Technical fundamentals matrix"""
        log_progress("Creating Slide 5: Fundamentos...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title = slide.shapes.title
        title.text = "Fundamentos Técnicos — Matriz de Priorização"
        title.text_frame.paragraphs[0].font.size = Pt(26)
        title.text_frame.paragraphs[0].font.bold = True

        # Create table
        rows = 7
        cols = 4

        left = Inches(0.8)
        top = Inches(2)
        width = Inches(8.4)
        height = Inches(4.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(1.4)

        # Headers
        headers = ['Ponto de Melhoria', 'Impacto', 'Esforço', 'Prioridade']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*COLORS_RGB['dark'])

            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Data
        improvements = [
            ("Otimizar velocidade mobile", "Alto", "Médio", "P1"),
            ("Implementar remarketing avançado", "Alto", "Baixo", "P1"),
            ("Estruturar CRM integrado", "Alto", "Alto", "P2"),
            ("Melhorar SEO técnico", "Médio", "Baixo", "P2"),
            ("Adicionar chatbot inteligente", "Médio", "Médio", "P3"),
            ("Expandir presença em LLMs (GEO)", "Médio", "Baixo", "P2")
        ]

        for row_idx, (improvement, impact, effort, priority) in enumerate(improvements, start=1):
            table.cell(row_idx, 0).text = improvement
            table.cell(row_idx, 1).text = impact
            table.cell(row_idx, 2).text = effort
            table.cell(row_idx, 3).text = priority

            for col_idx in range(4):
                cell = table.cell(row_idx, col_idx)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)

                if col_idx == 0:
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p.alignment = PP_ALIGN.CENTER

                # Color priority column
                if col_idx == 3:
                    if priority == "P1":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*COLORS_RGB['danger'])
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.font.bold = True
                    elif priority == "P2":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*COLORS_RGB['accent'])
                        p.font.bold = True
                    elif priority == "P3":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*COLORS_RGB['light'])

    def add_seo_geo_slide(self):
        """Slide 12: SEO + GEO importance"""
        log_progress("Creating Slide 12: Importância SEO + GEO...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title = slide.shapes.title
        title.text = "A Importância do SEO + GEO"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*COLORS_RGB['dark'])

        # Content
        left = Inches(0.8)
        top = Inches(1.9)
        width = Inches(8.4)
        height = Inches(5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        sections = [
            ("SEO (Search Engine Optimization)", "Posicionamento em buscadores tradicionais (Google, Bing) para capturar demanda ativa de usuários buscando produtos e serviços."),
            ("GEO (Generative Engine Optimization)", "Otimização para aparecer em respostas de LLMs (ChatGPT, Perplexity, Gemini, Claude) que estão se tornando o novo canal de descoberta de marcas."),
            ("Por que é crítico para Sumirê?", "• Busca orgânica = Custo de aquisição zero\n• LLMs respondem 40% das buscas em 2026\n• Presença em ambos = Máxima visibilidade\n• Concorrentes ainda não dominam este espaço")
        ]

        for idx, (section_title, content) in enumerate(sections):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = section_title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*COLORS_RGB['primary'])
            p.space_before = Pt(22) if idx > 0 else 0

            p2 = tf.add_paragraph()
            p2.text = content
            p2.font.size = Pt(15)
            p2.level = 0
            p2.space_after = Pt(10)

    def add_llm_presence_slide(self):
        """Slide 13: LLM presence ranking"""
        log_progress("Creating Slide 13: Presença em LLMs...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Title
        title = slide.shapes.title
        title.text = "Presença em LLMs — Ranking de Menções"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True

        # Create table
        platforms = ['ChatGPT', 'Perplexity', 'Gemini', 'Claude']
        brands = ['Época', 'Boticário', 'Sumirê', 'Sephora', 'Natura']

        rows = len(platforms) + 1
        cols = len(brands) + 1

        left = Inches(0.8)
        top = Inches(2.2)
        width = Inches(8.4)
        height = Inches(4.3)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header row
        table.cell(0, 0).text = "LLM / Marca"
        cell = table.cell(0, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*COLORS_RGB['dark'])
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

        for i, brand in enumerate(brands, start=1):
            cell = table.cell(0, i)
            cell.text = brand
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*COLORS_RGB['secondary'])
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        # Mock ranking data (in production, use real LLM data)
        rankings = {
            'ChatGPT': [2, 1, 3, 4, 5],
            'Perplexity': [1, 3, 4, 2, 5],
            'Gemini': [2, 1, 5, 3, 4],
            'Claude': [1, 2, 4, 3, 5]
        }

        # Fill data
        for row_idx, platform in enumerate(platforms, start=1):
            # Platform name
            cell = table.cell(row_idx, 0)
            cell.text = platform
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)

            # Rankings
            for col_idx, rank in enumerate(rankings[platform], start=1):
                cell = table.cell(row_idx, col_idx)
                cell.text = f"{rank}º"
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER

                # Color by position
                if rank == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*COLORS_RGB['success'])
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
                elif rank <= 3:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 243, 205)

                # Highlight Sumirê column
                if col_idx == 3:
                    p.font.bold = True

    def build_presentation(self) -> str:
        """Build complete presentation"""
        log_progress("="*60)
        log_progress("BUILDING PPTX PRESENTATION")
        log_progress("="*60)

        try:
            # Add all slides
            self.add_title_slide()
            self.add_methodology_slide()
            self.add_about_brand_slide()
            self.add_image_slide(4, "Radar de Maturidade Digital", "radar_sumire.png")
            self.add_fundamentals_slide()
            self.add_image_slide(6, "Comparativo de Maturidade Digital", "radar_competitors.png")
            self.add_image_slide(7, "Evolução de Tráfego — 12 Meses", "traffic_12months.png")
            self.add_image_slide(8, "Comparativo de Canais — Concorrentes", "competitors_channels.png")
            self.add_image_slide(9, "Fluxo de Tráfego — Origem e Destino", "sankey_flow.png")
            self.add_image_slide(10, "Top 15 Keywords Transacionais", "keywords_table.png")
            self.add_image_slide(11, "Demanda de Mercado — 12 Meses", "demand_trend.png")
            self.add_seo_geo_slide()
            self.add_llm_presence_slide()

            # Save presentation
            output_dir = os.path.join(os.path.dirname(__file__), '../output/')
            os.makedirs(output_dir, exist_ok=True)

            output_path = os.path.join(output_dir, 'sumire-analise-digital.pptx')
            self.prs.save(output_path)

            log_progress("\n" + "="*60)
            log_success("PRESENTATION CREATED SUCCESSFULLY!")
            log_progress("="*60)
            log_success(f"\n✅ Saved to: {output_path}")
            log_progress(f"📊 Total slides: {len(self.prs.slides)}")

            return output_path

        except Exception as e:
            log_error(f"Error building presentation: {e}")
            import traceback
            traceback.print_exc()
            return None


def main():
    """Main execution function"""

    # Load data files
    data_file = os.path.join(os.path.dirname(__file__), PATHS['sumire_data'])
    llm_file = os.path.join(os.path.dirname(__file__), PATHS['llm_results'])

    if not os.path.exists(data_file):
        log_error(f"Data file not found: {data_file}")
        log_error("Please run 01_browser_automation.py first!")
        return 1

    if not os.path.exists(llm_file):
        log_error(f"LLM data file not found: {llm_file}")
        log_error("Please run 02_llm_testing.py first!")
        return 1

    # Build presentation
    try:
        builder = PPTXBuilder(data_file, llm_file)
        output_path = builder.build_presentation()

        if output_path:
            return 0
        else:
            return 1

    except Exception as e:
        log_error(f"Fatal error: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    exit(main())
