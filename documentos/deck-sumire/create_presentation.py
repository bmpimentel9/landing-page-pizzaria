#!/usr/bin/env python3
"""
Perfumaria Sumire - Proposta Inicial v4 Company
Professional McKinsey/Accenture-style Consulting Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os
import json

# =============================================================================
# CONFIGURATION & COLOR SCHEME
# =============================================================================

# V4 Company Brand Colors (Professional Consulting Palette)
COLORS = {
    'primary_dark': RGBColor(31, 71, 136),      # #1F4788 - Deep Navy
    'primary_medium': RGBColor(68, 114, 196),   # #4472C4 - Medium Blue
    'primary_light': RGBColor(142, 169, 219),   # #8EA9DB - Light Blue
    'accent_orange': RGBColor(255, 107, 53),    # #FF6B35 - V4 Orange
    'accent_green': RGBColor(112, 173, 71),     # #70AD47 - Success Green
    'accent_red': RGBColor(192, 80, 77),        # #C0504D - Alert Red
    'text_dark': RGBColor(51, 51, 51),          # #333333 - Dark Gray
    'text_medium': RGBColor(89, 89, 89),        # #595959 - Medium Gray
    'text_light': RGBColor(127, 127, 127),      # #7F7F7F - Light Gray
    'background_light': RGBColor(242, 242, 242), # #F2F2F2 - Light Gray BG
    'white': RGBColor(255, 255, 255),           # #FFFFFF - White
}

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_shape_fill(shape, color):
    """Set shape fill color"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_text_frame_properties(text_frame, margin_left=Inches(0.1), margin_top=Inches(0.05)):
    """Set text frame properties"""
    text_frame.word_wrap = True
    text_frame.margin_left = margin_left
    text_frame.margin_right = margin_left
    text_frame.margin_top = margin_top
    text_frame.margin_bottom = margin_top

def add_text_to_shape(shape, text, font_name='Calibri', font_size=12,
                      bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add formatted text to a shape"""
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return text_frame

def add_title_shape(slide, text, left, top, width, height,
                    font_size=28, bold=True, color=None):
    """Add a title text box"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = add_text_to_shape(shape, text, font_size=font_size, bold=bold,
                           color=color or COLORS['text_dark'])
    return shape

def add_rectangle(slide, left, top, width, height, fill_color):
    """Add a rectangle shape with fill"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def format_number(num):
    """Format number with thousand separators (Brazilian format)"""
    return f"{num:,.0f}".replace(",", ".")

# =============================================================================
# SLIDE CREATION FUNCTIONS
# =============================================================================

def create_cover_slide(prs):
    """Slide 1: Cover/Capa - Professional title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Top accent bar
    top_bar = add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), COLORS['accent_orange'])

    # Main background area (left side - dark blue)
    left_bg = add_rectangle(slide, Inches(0), Inches(0.15), Inches(5), Inches(7.35), COLORS['primary_dark'])

    # V4 Company Logo placeholder (text-based)
    logo_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.8))
    tf = logo_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "v4"
    run.font.name = 'Calibri'
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = COLORS['accent_orange']
    run2 = p.add_run()
    run2.text = " company"
    run2.font.name = 'Calibri'
    run2.font.size = Pt(28)
    run2.font.color.rgb = COLORS['white']

    # Main title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4.3), Inches(1.5))
    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Analise Digital"
    run.font.name = 'Calibri'
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    # Subtitle
    subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(4.3), Inches(0.8))
    tf = subtitle_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Oportunidades de Transformacao Digital"
    run.font.name = 'Calibri'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['primary_light']

    # Client name - prominent
    client_shape = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(7.3), Inches(1.2))
    tf = client_shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Perfumaria Sumire"
    run.font.name = 'Calibri'
    run.font.size = Pt(52)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    # Client tagline
    tagline_shape = slide.shapes.add_textbox(Inches(5.5), Inches(3.7), Inches(7.3), Inches(0.6))
    tf = tagline_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Maior rede de perfumarias do Brasil"
    run.font.name = 'Calibri'
    run.font.size = Pt(22)
    run.font.italic = True
    run.font.color.rgb = COLORS['text_medium']

    # Key stats
    stats_data = [
        ("40+", "Anos de Historia"),
        ("~70", "Lojas em SP"),
        ("124K", "Seguidores Instagram")
    ]

    for i, (number, label) in enumerate(stats_data):
        x_pos = Inches(5.5 + i * 2.4)

        # Number
        num_shape = slide.shapes.add_textbox(x_pos, Inches(4.8), Inches(2.2), Inches(0.8))
        tf = num_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = number
        run.font.name = 'Calibri'
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = COLORS['accent_orange']

        # Label
        label_shape = slide.shapes.add_textbox(x_pos, Inches(5.4), Inches(2.2), Inches(0.5))
        tf = label_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = 'Calibri'
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['text_medium']

    # Date and meeting info
    date_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(4), Inches(0.5))
    tf = date_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Janeiro 2026"
    run.font.name = 'Calibri'
    run.font.size = Pt(14)
    run.font.color.rgb = COLORS['primary_light']

    # Confidential note
    conf_shape = slide.shapes.add_textbox(Inches(5.5), Inches(6.8), Inches(7), Inches(0.4))
    tf = conf_shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "CONFIDENCIAL | Proposta Comercial Inicial"
    run.font.name = 'Calibri'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['text_light']

    return slide

def create_agenda_slide(prs):
    """Slide 2: Agenda - Meeting flow"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header_bar = add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), COLORS['primary_dark'])

    # Slide title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(8), Inches(0.7))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Agenda da Reuniao"
    run.font.name = 'Calibri'
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    # Subtitle
    subtitle_shape = slide.shapes.add_textbox(Inches(9), Inches(0.4), Inches(4), Inches(0.6))
    tf = subtitle_shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "35 minutos"
    run.font.name = 'Calibri'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['accent_orange']

    # Agenda items
    agenda_items = [
        ("01", "Visao Geral do Cliente", "5 min", "Contexto de mercado e posicionamento atual"),
        ("02", "Analise de Audiencia Digital", "8 min", "Trafego, canais e comportamento do consumidor"),
        ("03", "Radar de Maturidade Digital", "7 min", "Avaliacao comparativa vs. concorrentes"),
        ("04", "Gaps e Oportunidades", "8 min", "Pontos de melhoria com impacto de negocio"),
        ("05", "Keywords Transacionais", "5 min", "Oportunidades de busca organica e paga"),
        ("06", "Proximos Passos", "2 min", "Recomendacoes e proposta de parceria"),
    ]

    for i, (num, title, time, desc) in enumerate(agenda_items):
        y_pos = Inches(1.5 + i * 0.95)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y_pos, Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['accent_orange']
        circle.line.fill.background()

        # Number text
        num_shape = slide.shapes.add_textbox(Inches(0.5), y_pos + Inches(0.1), Inches(0.6), Inches(0.5))
        tf = num_shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.name = 'Calibri'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = COLORS['white']

        # Title
        title_shape = slide.shapes.add_textbox(Inches(1.3), y_pos + Inches(0.05), Inches(6), Inches(0.4))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = 'Calibri'
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = COLORS['text_dark']

        # Description
        desc_shape = slide.shapes.add_textbox(Inches(1.3), y_pos + Inches(0.4), Inches(7), Inches(0.35))
        tf = desc_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.name = 'Calibri'
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['text_medium']

        # Time
        time_shape = slide.shapes.add_textbox(Inches(11.5), y_pos + Inches(0.15), Inches(1.5), Inches(0.4))
        tf = time_shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = time
        run.font.name = 'Calibri'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS['primary_medium']

    # Bottom note
    note_shape = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.4))
    tf = note_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "* Dados obtidos atraves da tecnologia proprietaria v4 Company - Inteligencia Digital 2026"
    run.font.name = 'Calibri'
    run.font.size = Pt(10)
    run.font.italic = True
    run.font.color.rgb = COLORS['text_light']

    return slide

def create_client_overview_slide(prs, data):
    """Slide 3: Client Overview"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header_bar = add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), COLORS['primary_dark'])

    # Slide title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.7))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Perfumaria Sumire: Lider Regional com Oportunidade Digital"
    run.font.name = 'Calibri'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    # Company description box
    desc_box = add_rectangle(slide, Inches(0.5), Inches(1.5), Inches(7.5), Inches(2.5), COLORS['background_light'])

    desc_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(7.1), Inches(2.2))
    tf = desc_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Sobre a Empresa"
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "\nFundada em 1984, a Perfumaria Sumire se consolidou como a maior rede de perfumarias do Brasil, com aproximadamente 70 lojas no estado de Sao Paulo. A empresa combina tradicao familiar com presenca omnichannel, operando e-commerce integrado com varejo fisico."
    run2.font.name = 'Calibri'
    run2.font.size = Pt(13)
    run2.font.color.rgb = COLORS['text_dark']

    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = "\nPosicionamento: Perfumes importados e nacionais com atendimento especializado e precos competitivos."
    run3.font.name = 'Calibri'
    run3.font.size = Pt(13)
    run3.font.bold = True
    run3.font.color.rgb = COLORS['text_medium']

    # Key metrics cards
    metrics = [
        ("2.04M", "Visitas Anuais", "Website 2025"),
        ("~70", "Lojas Fisicas", "Estado de SP"),
        ("124K", "Seguidores", "Instagram"),
        ("1.87%", "Engajamento", "Redes Sociais"),
    ]

    card_width = Inches(1.85)
    start_x = Inches(8.3)

    for i, (value, label, sublabel) in enumerate(metrics):
        y_offset = Inches(1.5 + i * 1.5)

        # Card background
        card = add_rectangle(slide, start_x, y_offset, card_width, Inches(1.3), COLORS['white'])
        card.shadow.inherit = False

        # Left accent
        accent = add_rectangle(slide, start_x, y_offset, Inches(0.08), Inches(1.3), COLORS['accent_orange'])

        # Value
        val_shape = slide.shapes.add_textbox(start_x + Inches(0.15), y_offset + Inches(0.15), card_width - Inches(0.2), Inches(0.5))
        tf = val_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.name = 'Calibri'
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = COLORS['primary_dark']

        # Label
        lbl_shape = slide.shapes.add_textbox(start_x + Inches(0.15), y_offset + Inches(0.6), card_width - Inches(0.2), Inches(0.35))
        tf = lbl_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = 'Calibri'
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = COLORS['text_dark']

        # Sublabel
        sub_shape = slide.shapes.add_textbox(start_x + Inches(0.15), y_offset + Inches(0.9), card_width - Inches(0.2), Inches(0.3))
        tf = sub_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = sublabel
        run.font.name = 'Calibri'
        run.font.size = Pt(10)
        run.font.color.rgb = COLORS['text_light']

    # Technical assessment section
    tech_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4), Inches(0.4))
    tf = tech_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Avaliacao Tecnica Atual"
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    # Tech metrics
    tech_data = data['sumire']['technical']
    tech_items = [
        ("PageSpeed Mobile", f"{tech_data['pagespeed_mobile']}/100", tech_data['pagespeed_mobile'] >= 75),
        ("PageSpeed Desktop", f"{tech_data['pagespeed_desktop']}/100", tech_data['pagespeed_desktop'] >= 80),
        ("SEO Score", f"{tech_data['seo_score']}/100", tech_data['seo_score'] >= 75),
        ("Acessibilidade", f"{tech_data['accessibility_score']}/100", tech_data['accessibility_score'] >= 75),
        ("Google Analytics", "Implementado" if tech_data['google_analytics'] else "Ausente", tech_data['google_analytics']),
        ("Facebook Pixel", "Implementado" if tech_data['facebook_pixel'] else "Ausente", tech_data['facebook_pixel']),
        ("Google Tag Manager", "Implementado" if tech_data['google_tag_manager'] else "Ausente", tech_data['google_tag_manager']),
    ]

    for i, (metric, value, is_good) in enumerate(tech_items):
        col = i % 4
        row = i // 4
        x_pos = Inches(0.5 + col * 2.0)
        y_pos = Inches(4.7 + row * 0.9)

        # Metric name
        m_shape = slide.shapes.add_textbox(x_pos, y_pos, Inches(1.9), Inches(0.4))
        tf = m_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = metric
        run.font.name = 'Calibri'
        run.font.size = Pt(11)
        run.font.color.rgb = COLORS['text_medium']

        # Value with color
        v_shape = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.35), Inches(1.9), Inches(0.4))
        tf = v_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.name = 'Calibri'
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLORS['accent_green'] if is_good else COLORS['accent_red']

    # Key insight box
    insight_box = add_rectangle(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8), COLORS['primary_light'])

    insight_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(12), Inches(0.6))
    tf = insight_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "INSIGHT: "
    run.font.name = 'Calibri'
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']
    run2 = p.add_run()
    run2.text = "Ausencia de Google Tag Manager limita a capacidade de rastreamento avancado e integracao com plataformas de automacao de marketing."
    run2.font.name = 'Calibri'
    run2.font.size = Pt(12)
    run2.font.color.rgb = COLORS['text_dark']

    return slide

def create_audience_radar_slide(prs, data):
    """Slide 4: Audience Analysis & Maturity Radar with embedded charts"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header_bar = add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), COLORS['primary_dark'])

    # Slide title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Sumire vs. Mercado: Trafego Solido, Maturidade Digital a Desenvolver"
    run.font.name = 'Calibri'
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    # Left section - Traffic Analysis
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Evolucao de Trafego 12 Meses"
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    # Insert traffic chart
    chart_path = '/Users/brunopimentelV4/Desktop/deck-sumire/charts/traffic_12months.png'
    if os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(0.3), Inches(1.9), width=Inches(6.2))

    # Right section - Radar
    right_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Radar de Maturidade Digital"
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    # Insert radar chart
    radar_path = '/Users/brunopimentelV4/Desktop/deck-sumire/charts/radar_sumire.png'
    if os.path.exists(radar_path):
        slide.shapes.add_picture(radar_path, Inches(6.6), Inches(1.8), width=Inches(6.2))

    # Key metrics summary
    total_traffic_2025 = sum(data['sumire']['traffic_12m']['total'])

    metrics_box = add_rectangle(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.9), COLORS['background_light'])

    summary_metrics = [
        (f"{format_number(total_traffic_2025)}", "Visitas em 2025", COLORS['primary_dark']),
        ("+92%", "Crescimento YoY", COLORS['accent_green']),
        ("40%", "Trafego Direto", COLORS['primary_medium']),
        ("28%", "Busca Organica", COLORS['primary_medium']),
        ("15%", "Redes Sociais", COLORS['accent_orange']),
    ]

    for i, (val, label, color) in enumerate(summary_metrics):
        x_pos = Inches(0.8 + i * 2.5)

        val_shape = slide.shapes.add_textbox(x_pos, Inches(6.5), Inches(2.3), Inches(0.45))
        tf = val_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = val
        run.font.name = 'Calibri'
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = color

        lbl_shape = slide.shapes.add_textbox(x_pos, Inches(6.9), Inches(2.3), Inches(0.3))
        tf = lbl_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = 'Calibri'
        run.font.size = Pt(10)
        run.font.color.rgb = COLORS['text_medium']

    return slide

def create_gap_analysis_slide(prs, data, llm_data):
    """Slide 5: Maturity Gap Analysis"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header_bar = add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), COLORS['primary_dark'])

    # Slide title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Gaps Identificados: Oportunidades de Alto Impacto"
    run.font.name = 'Calibri'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    # Gap cards with detailed explanations
    gaps = [
        {
            "title": "Ausencia de Google Tag Manager",
            "impact": "ALTO",
            "description": "Sem GTM, cada alteracao de tracking exige desenvolvimento. Isso impede testes A/B rapidos, integracao com CRMs e automacao de marketing. Concorrentes com GTM implementam campanhas 60% mais rapido.",
            "recommendation": "Implementar GTM com camada de dados estruturada para e-commerce"
        },
        {
            "title": "Score de Acessibilidade Baixo (70/100)",
            "impact": "MEDIO",
            "description": "Problemas de contraste, navegacao por teclado e leitores de tela. Alem de excluir 15% do publico potencial, afeta negativamente o ranking no Google, que prioriza sites acessiveis.",
            "recommendation": "Auditoria WCAG 2.1 e correcoes prioritarias de navegacao"
        },
        {
            "title": "Presenca Limitada em LLMs/IA",
            "impact": "ALTO",
            "description": f"Sumire aparece em apenas 57% das buscas relevantes em ChatGPT, Perplexity, Gemini e Claude. Posicao media: 3.2 (vs. Soneda: 2.8). A nova geracao de consumidores usa IA para descoberta de marcas.",
            "recommendation": "Estrategia de conteudo otimizada para citacao em LLMs"
        },
        {
            "title": "Trafego Pago Subotimizado (5%)",
            "impact": "ALTO",
            "description": "Apenas 5% do trafego vem de midia paga, enquanto Lojas REDE investe ~6% e Epoca 6%. Ha oportunidade de capturar demanda transacional com ROI mensuravel em keywords de alta intencao.",
            "recommendation": "Campanha piloto Google Ads focada em top 10 keywords transacionais"
        },
    ]

    for i, gap in enumerate(gaps):
        col = i % 2
        row = i // 2
        x_pos = Inches(0.4 + col * 6.4)
        y_pos = Inches(1.4 + row * 2.9)

        # Card background
        card = add_rectangle(slide, x_pos, y_pos, Inches(6.2), Inches(2.7), COLORS['white'])

        # Impact indicator
        impact_color = COLORS['accent_red'] if gap['impact'] == 'ALTO' else COLORS['accent_orange']
        impact_bar = add_rectangle(slide, x_pos, y_pos, Inches(0.12), Inches(2.7), impact_color)

        # Impact badge
        badge_shape = slide.shapes.add_textbox(x_pos + Inches(5.0), y_pos + Inches(0.1), Inches(1), Inches(0.3))
        tf = badge_shape.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"IMPACTO {gap['impact']}"
        run.font.name = 'Calibri'
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = impact_color

        # Title
        title_shape = slide.shapes.add_textbox(x_pos + Inches(0.25), y_pos + Inches(0.15), Inches(5.5), Inches(0.45))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = gap['title']
        run.font.name = 'Calibri'
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLORS['text_dark']

        # Description
        desc_shape = slide.shapes.add_textbox(x_pos + Inches(0.25), y_pos + Inches(0.6), Inches(5.8), Inches(1.3))
        tf = desc_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = gap['description']
        run.font.name = 'Calibri'
        run.font.size = Pt(11)
        run.font.color.rgb = COLORS['text_medium']

        # Recommendation
        rec_shape = slide.shapes.add_textbox(x_pos + Inches(0.25), y_pos + Inches(2.0), Inches(5.8), Inches(0.6))
        tf = rec_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "-> "
        run.font.name = 'Calibri'
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = COLORS['accent_green']
        run2 = p.add_run()
        run2.text = gap['recommendation']
        run2.font.name = 'Calibri'
        run2.font.size = Pt(10)
        run2.font.color.rgb = COLORS['accent_green']

    # Bottom summary
    summary_bar = add_rectangle(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35), COLORS['primary_light'])

    summary_text = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12), Inches(0.3))
    tf = summary_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Potencial de crescimento estimado: +35% em trafego qualificado e +25% em conversao com implementacao das recomendacoes"
    run.font.name = 'Calibri'
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    return slide

def create_keywords_slide(prs, data):
    """Slide 6: Transactional Keywords Analysis with embedded chart"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header_bar = add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), COLORS['primary_dark'])

    # Slide title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Keywords Transacionais: 102K Buscas Mensais em Oportunidade"
    run.font.name = 'Calibri'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    # Insert keywords table chart
    keywords_chart_path = '/Users/brunopimentelV4/Desktop/deck-sumire/charts/keywords_table.png'
    if os.path.exists(keywords_chart_path):
        slide.shapes.add_picture(keywords_chart_path, Inches(0.3), Inches(1.4), width=Inches(8.5))

    # Right side - Key insights
    insights_title = slide.shapes.add_textbox(Inches(9), Inches(1.4), Inches(4), Inches(0.4))
    tf = insights_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Destaques"
    run.font.name = 'Calibri'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    insights = [
        ("102K+", "Volume mensal total"),
        ("R$2.20", "CPC medio"),
        ("Pos. 2", "Melhor ranking atual"),
        ("6 de 15", "Keywords no Top 10"),
    ]

    for i, (val, label) in enumerate(insights):
        y_pos = Inches(1.9 + i * 0.9)

        # Card
        card = add_rectangle(slide, Inches(9), y_pos, Inches(4), Inches(0.75), COLORS['background_light'])

        val_shape = slide.shapes.add_textbox(Inches(9.2), y_pos + Inches(0.1), Inches(1.5), Inches(0.4))
        tf = val_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = val
        run.font.name = 'Calibri'
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = COLORS['accent_orange']

        lbl_shape = slide.shapes.add_textbox(Inches(10.7), y_pos + Inches(0.2), Inches(2.1), Inches(0.4))
        tf = lbl_shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.name = 'Calibri'
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['text_dark']

    # Opportunity analysis box
    opp_title = slide.shapes.add_textbox(Inches(9), Inches(5.5), Inches(4), Inches(0.35))
    tf = opp_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Oportunidade Imediata"
    run.font.name = 'Calibri'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['primary_dark']

    opp_box = add_rectangle(slide, Inches(9), Inches(5.9), Inches(4), Inches(1.3), COLORS['primary_light'])

    opp_text = slide.shapes.add_textbox(Inches(9.2), Inches(6.0), Inches(3.6), Inches(1.1))
    tf = opp_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "9 keywords com posicao >5 e volume >5K representam oportunidade de R$180K/ano em trafego organico equivalente (vs. custo de aquisicao paga)."
    run.font.name = 'Calibri'
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS['text_dark']

    # Bottom CTA
    cta_bar = add_rectangle(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.35), COLORS['accent_orange'])

    cta_text = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12), Inches(0.3))
    tf = cta_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Proximo passo: Proposta detalhada de SEO + Paid Search com projecao de ROI em 90 dias"
    run.font.name = 'Calibri'
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = COLORS['white']

    return slide

# =============================================================================
# MAIN EXECUTION
# =============================================================================

def main():
    """Create the complete presentation"""

    # Load data
    with open('/Users/brunopimentelV4/Desktop/deck-sumire/data/sumire_data.json', 'r') as f:
        data = json.load(f)

    with open('/Users/brunopimentelV4/Desktop/deck-sumire/data/llm_results.json', 'r') as f:
        llm_data = json.load(f)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Creating Perfumaria Sumire - v4 Company Consulting Presentation...")
    print("=" * 60)

    # Create slides
    print("  [1/6] Creating Cover Slide...")
    create_cover_slide(prs)

    print("  [2/6] Creating Agenda Slide...")
    create_agenda_slide(prs)

    print("  [3/6] Creating Client Overview Slide...")
    create_client_overview_slide(prs, data)

    print("  [4/6] Creating Audience & Radar Slide...")
    create_audience_radar_slide(prs, data)

    print("  [5/6] Creating Gap Analysis Slide...")
    create_gap_analysis_slide(prs, data, llm_data)

    print("  [6/6] Creating Keywords Analysis Slide...")
    create_keywords_slide(prs, data)

    # Save presentation
    output_path = '/Users/brunopimentelV4/Desktop/deck-sumire/output/Perfumaria_Sumire_Proposta_Inicial_v4company.pptx'
    prs.save(output_path)

    print("=" * 60)
    print(f"SUCCESS: Presentation saved to:")
    print(f"  {output_path}")
    print("=" * 60)

    return output_path

if __name__ == "__main__":
    main()
